#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 课件生成脚本：把 data/CKGDT/课件/*.md 转成 .pptx 课件。

依赖：python-pptx

用法：
    python tools/generate_ckgdt_pptx.py            # 生成所有课件
    python tools/generate_ckgdt_pptx.py 第一章      # 只生成某章课件
    python tools/generate_ckgdt_pptx.py --list      # 列出所有课件大纲
"""

import os
import re
import sys
from pathlib import Path

# Windows 终端 UTF-8
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

ROOT = Path(__file__).resolve().parent.parent
SRC_DIR = ROOT / "data" / "CKGDT" / "课件"
OUT_DIR = ROOT / "data" / "CKGDT" / "输出" / "课件"

# PPTX 配置
SLIDE_WIDTH = 12192000  # 33.87cm (16:9)
SLIDE_HEIGHT = 6858000  # 19.05cm
FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
COLOR_BG = (15, 23, 42)  # 深蓝 #0f172a
COLOR_ACCENT = (102, 126, 234)  # 主题色 #667eea
COLOR_WHITE = (255, 255, 255)
COLOR_GRAY = (148, 163, 184)  # slate-400


def parse_outline(md_path: Path) -> dict:
    """解析课件大纲 Markdown。"""
    content = md_path.read_text(encoding="utf-8")
    
    # 提取标题
    title_match = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
    title = title_match.group(1).strip() if title_match else md_path.stem
    
    # 提取各节
    sections = []
    current_section = None
    
    for line in content.split("\n"):
        line = line.strip()
        if line.startswith("## "):
            if current_section:
                sections.append(current_section)
            current_section = {
                "title": line[3:].strip(),
                "bullets": [],
                "notes": "",
            }
        elif line.startswith("- ") and current_section:
            current_section["bullets"].append(line[2:].strip())
        elif line.startswith("> ") and current_section:
            current_section["notes"] += line[2:].strip() + "\n"
    
    if current_section:
        sections.append(current_section)
    
    return {
        "title": title,
        "sections": sections,
    }


def create_slide(presentation, title: str, bullets: list, notes: str = "",
                 is_title_slide: bool = False):
    """添加一张幻灯片。"""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    
    slide_layout = presentation.slide_layouts[6]  # 空白布局
    slide = presentation.slides.add_slide(slide_layout)
    
    # 设置深色背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*COLOR_BG)
    
    if is_title_slide:
        # 标题页
        left = Inches(1)
        top = Inches(3)
        width = Inches(11)
        height = Inches(2)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLOR_ACCENT)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if bullets:
            p2 = tf.add_paragraph()
            p2.text = bullets[0]
            p2.font.name = FONT_BODY
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(*COLOR_GRAY)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(20)
    else:
        # 内容页
        # 标题
        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(11.4)
        height = Inches(1)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLOR_WHITE)
        
        # 要点
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(11)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = FONT_BODY
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(*COLOR_WHITE)
            p.space_before = Pt(12)
            p.level = 0
        
        # 备注
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
    
    return slide


def generate_pptx(outline: dict, output_dir: Path) -> Path:
    """根据大纲生成 PPTX。"""
    from pptx import Presentation
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    title = outline["title"]
    sections = outline["sections"]
    
    # 标题页
    create_slide(prs, title, [f"课程知识图谱与数字孪生平台"], is_title_slide=True)
    
    # 目录页
    toc_bullets = [s["title"] for s in sections]
    create_slide(prs, "目录", toc_bullets)
    
    # 各节内容
    for section in sections:
        create_slide(prs, section["title"], section["bullets"], section["notes"])
    
    # 结束页
    create_slide(prs, "谢谢", ["课程知识图谱与数字孪生平台"], is_title_slide=True)
    
    # 保存
    output_path = output_dir / f"{title}.pptx"
    prs.save(str(output_path))
    print(f"  [OK] {output_path.name}")
    return output_path


def list_outlines():
    """列出所有课件大纲。"""
    outlines = sorted(SRC_DIR.glob("*.md"))
    print(f"\n共 {len(outlines)} 个课件大纲：")
    for o in outlines:
        print(f"  - {o.name}")


def main():
    if "--list" in sys.argv:
        list_outlines()
        return
    
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    
    # 确定要生成的大纲
    if len(sys.argv) > 1 and sys.argv[1] != "--list":
        keyword = sys.argv[1]
        outlines = [s for s in SRC_DIR.glob("*.md") if keyword in s.name]
        if not outlines:
            print(f"[ERROR] 未找到包含 '{keyword}' 的大纲")
            return
    else:
        outlines = sorted(SRC_DIR.glob("*.md"))
    
    print(f"\n=== CKGDT PPTX 课件生成 ===")
    print(f"大纲目录: {SRC_DIR}")
    print(f"输出目录: {OUT_DIR}")
    print(f"待处理: {len(outlines)} 个大纲\n")
    
    success = 0
    for outline_path in outlines:
        print(f"[{outline_path.name}]")
        try:
            outline = parse_outline(outline_path)
            result = generate_pptx(outline, OUT_DIR)
            if result:
                success += 1
        except Exception as e:
            print(f"  [ERROR] {e}")
    
    print(f"\n完成: {success}/{len(outlines)} 个课件")


if __name__ == "__main__":
    main()
