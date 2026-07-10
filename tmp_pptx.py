
import json, sys, os

def main():
    if len(sys.argv) < 2:
        print("Usage: python gen_pptx.py data.json", file=sys.stderr)
        sys.exit(1)

    with open(sys.argv[1], 'r', encoding='utf-8') as f:
        data = json.load(f)

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    CJK_FONT = '\u5fae\u8f6f\u96c5\u9ed1'

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = data.get('title', '\u6559\u5b66\u8bfe\u4ef6')
    chapter = data.get('chapter', '')
    slides_data = data.get('slides', [])
    output = data.get('output', 'output.pptx')

    # ── 色彩主题 ──
    PRIMARY = RGBColor(0x0D, 0x47, 0xA1)       # 深蓝
    SECONDARY = RGBColor(0x1E, 0x88, 0xE5)      # 中蓝
    ACCENT = RGBColor(0x00, 0xAC, 0xC1)          # 青色
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
    DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
    MID_TEXT = RGBColor(0x4A, 0x4A, 0x5A)
    LIGHT_TEXT = RGBColor(0x8A, 0x8A, 0x9A)
    CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
    CODE_FG = RGBColor(0xCD, 0xD6, 0xF4)
    SECTION_BG = RGBColor(0xE3, 0xF2, 0xFD)
    CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
    ORANGE = RGBColor(0xFF, 0x6F, 0x00)

    def set_font(p, size, color, bold=False, name=CJK_FONT):
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = name
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = name
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {})
                rPr.append(ea)
            ea.set('typeface', name)

    def set_p_font(p, text, size, color, bold=False, name=CJK_FONT):
        p.text = text
        set_font(p, size, color, bold, name)

    def add_gradient_bg(slide, c1=PRIMARY, c2=SECONDARY):
        """Set 135-degree gradient background."""
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_angle = 135
        stops = fill.gradient_stops
        stops[0].color.rgb = c1
        stops[0].position = 0.0
        stops[1].color.rgb = c2
        stops[1].position = 1.0

    def add_solid_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape_with_gradient(slide, left, top, width, height, c1, c2, corner_radius=0):
        """Add a rounded rectangle with gradient fill."""
        shape = slide.shapes.add_shape(
            5, left, top, width, height  # 5 = rounded rectangle
        )
        shape.fill.gradient()
        shape.fill.gradient_angle = 135
        stops = shape.fill.gradient_stops
        stops[0].color.rgb = c1
        stops[0].position = 0.0
        stops[1].color.rgb = c2
        stops[1].position = 1.0
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def add_card(slide, left, top, width, height, shadow=True):
        """Add a white card with subtle shadow."""
        shape = slide.shapes.add_shape(
            5, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.width = Pt(0.5)
        if shadow:
            shape.shadow.inherit = False
        return shape

    def add_entrance_anim(slide, shape, delay_ms=0, effect='fade'):
        """Add entrance animation to a shape."""
        try:
            tree = slide._element
            timing = tree.find(qn('p:timing'))
            if timing is None:
                timing = etree.SubElement(tree, qn('p:timing'))
            tnLst = timing.find(qn('p:tnLst'))
            if tnLst is None:
                tnLst = etree.SubElement(timing, qn('p:tnLst'))

            par = tnLst.find(qn('p:par'))
            if par is None:
                par = etree.SubElement(tnLst, qn('p:par'))
                cTn_root = etree.SubElement(par, qn('p:cTn'), {
                    'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'
                })
                childTnLst = etree.SubElement(cTn_root, qn('p:childTnLst'))
                seq = etree.SubElement(childTnLst, qn('p:seq'), {'concurrent': '1', 'nextAc': 'seek'})
                seq_cTn = etree.SubElement(seq, qn('p:cTn'), {
                    'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'
                })
                etree.SubElement(seq_cTn, qn('p:childTnLst'))
                etree.SubElement(seq, qn('p:prevCondLst')).append(
                    etree.Element(qn('p:cond'), {'evt': 'onPrev', 'delay': '0'})
                )
                etree.SubElement(seq, qn('p:nextCondLst')).append(
                    etree.Element(qn('p:cond'), {'evt': 'onNext', 'delay': '0'})
                )
            else:
                cTn_root = par.find(qn('p:cTn'))
                childTnLst = cTn_root.find(qn('p:childTnLst'))
                seq = childTnLst.find(qn('p:seq'))
                seq_cTn = seq.find(qn('p:cTn'))

            main_childTnLst = seq_cTn.find(qn('p:childTnLst'))

            # Get max ID
            max_id = 2
            for el in tree.iter():
                ctn_id = el.get('id')
                if ctn_id and ctn_id.isdigit():
                    max_id = max(max_id, int(ctn_id))
            next_id = max_id + 1

            # par wrapper for this animation
            anim_par = etree.SubElement(main_childTnLst, qn('p:par'))
            anim_cTn = etree.SubElement(anim_par, qn('p:cTn'), {
                'id': str(next_id), 'fill': 'hold'
            })
            stCondLst = etree.SubElement(anim_cTn, qn('p:stCondLst'))
            etree.SubElement(stCondLst, qn('p:cond'), {'delay': str(delay_ms)})

            inner_childTnLst = etree.SubElement(anim_cTn, qn('p:childTnLst'))
            inner_par = etree.SubElement(inner_childTnLst, qn('p:par'))
            inner_cTn = etree.SubElement(inner_par, qn('p:cTn'), {
                'id': str(next_id + 1), 'presetID': '10' if effect == 'fade' else '2',
                'presetClass': 'entr', 'presetSubtype': '0',
                'fill': 'hold', 'nodeType': 'withEffect'
            })
            inner_stCond = etree.SubElement(inner_cTn, qn('p:stCondLst'))
            etree.SubElement(inner_stCond, qn('p:cond'), {'delay': '0'})

            inner_child = etree.SubElement(inner_cTn, qn('p:childTnLst'))

            # animEffect (fade)
            anim_effect = etree.SubElement(inner_child, qn('p:animEffect'), {
                'transition': 'in', 'filter': 'fade' if effect == 'fade' else 'wipe(down)'
            })
            eff_cBhvr = etree.SubElement(anim_effect, qn('p:cBhvr'))
            eff_cTn2 = etree.SubElement(eff_cBhvr, qn('p:cTn'), {
                'id': str(next_id + 2), 'dur': '500'
            })
            eff_tgtEl = etree.SubElement(eff_cBhvr, qn('p:tgtEl'))
            sp_tgt = etree.SubElement(eff_tgtEl, qn('p:spTgt'), {
                'spid': str(shape.shape_id)
            })

        except Exception as e:
            pass  # Animation is optional enhancement

    def add_slide_transition(slide, trans_type='fade'):
        """Add slide transition."""
        try:
            tree = slide._element
            transition = tree.find(qn('p:transition'))
            if transition is None:
                transition = etree.SubElement(tree, qn('p:transition'))
            transition.set('spd', 'med')
            transition.set('advClick', '1')
            if trans_type == 'fade':
                etree.SubElement(transition, qn('p:fade'))
            elif trans_type == 'push':
                etree.SubElement(transition, qn('p:push'))
            elif trans_type == 'wipe':
                etree.SubElement(transition, qn('p:wipe'))
        except:
            pass

    def parse_table_rows(rows):
        result = []
        for r in rows:
            cols = [c.strip() for c in r.strip().strip('|').split('|')]
            if cols and any(c for c in cols):
                result.append(cols)
        return result

    def add_table(slide, table_rows, left, top, width):
        try:
            parsed = parse_table_rows(table_rows)
            if not parsed:
                return top
            # Filter out separator rows (----)
            parsed = [r for r in parsed if not all(set(c.strip()) <= {'-', ':'} for c in r)]
            if not parsed:
                return top
            n_rows = len(parsed)
            n_cols = max(len(r) for r in parsed)
            for r in parsed:
                while len(r) < n_cols:
                    r.append('')
            row_h = Inches(0.42)
            col_w = Inches(width / n_cols)
            table_shape = slide.shapes.add_table(
                n_rows, n_cols,
                Inches(left), Inches(top),
                Inches(width), row_h * n_rows
            )
            tbl = table_shape.table
            for ci in range(n_cols):
                tbl.columns[ci].width = col_w
            for ri, row in enumerate(parsed):
                for ci, val in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = val
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if ri == 0:
                        set_font(p, 13, WHITE, bold=True)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = PRIMARY
                    else:
                        set_font(p, 12, DARK_TEXT)
                        if ri % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = SECTION_BG
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = LIGHT
                    cell.margin_left = Inches(0.08)
                    cell.margin_right = Inches(0.08)
                    cell.margin_top = Inches(0.04)
                    cell.margin_bottom = Inches(0.04)
            return top + (n_rows * 0.42) + 0.15
        except Exception as e:
            print(f"add_table error: {e}", file=sys.stderr)
            return top

    # ═══════════════════════════════════════════════════════════════
    # 封面页 — 全屏渐变 + 大标题 + 装饰线
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, PRIMARY, RGBColor(0x1A, 0x23, 0x7E))
    add_slide_transition(slide, 'fade')

    # 顶部装饰圆弧
    deco = slide.shapes.add_shape(
        9, Inches(-2), Inches(-3), Inches(17), Inches(6)  # 椭圆
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(0x1E, 0x88, 0xE5)
    deco.fill.fore_color.brightness = 0.05
    deco.line.fill.background()

    # 标题
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_p_font(p, title, 44, WHITE, bold=True)
    p.alignment = PP_ALIGN.CENTER
    add_entrance_anim(slide, tx, 200, 'fade')

    # 分隔线
    line = slide.shapes.add_shape(
        1, Inches(5), Inches(3.7), Inches(3.3), Emu(18000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if chapter:
        tx_ch = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.6))
        p_ch = tx_ch.text_frame.paragraphs[0]
        set_p_font(p_ch, chapter, 22, RGBColor(0xBB, 0xDD, 0xFF))
        p_ch.alignment = PP_ALIGN.CENTER
        add_entrance_anim(slide, tx_ch, 500, 'fade')

    # 底部信息
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.3), Inches(0.5))
    p3 = tx2.text_frame.paragraphs[0]
    set_p_font(p3, '\u79fb\u52a8\u5e94\u7528\u5f00\u53d1\u77e5\u8bc6\u56fe\u8c31\u6559\u5b66\u7cfb\u7edf', 14, RGBColor(0x88, 0x99, 0xCC))
    p3.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════════════════════════
    # 内容页 — 卡片式布局 + 渐变侧栏 + 动画入场
    # ═══════════════════════════════════════════════════════════════
    for idx, s in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_solid_bg(slide, LIGHT)
        add_slide_transition(slide, ['fade', 'push', 'wipe'][idx % 3])

        s_title = s.get('title', '')
        subtitle = s.get('subtitle', '')
        bullets = s.get('bullets', [])
        code = s.get('code', '')
        notes = s.get('notes', '')

        # ── 左侧渐变装饰条 ──
        side_bar = add_shape_with_gradient(
            slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height,
            PRIMARY, SECONDARY
        )

        # ── 顶部标题区域 ──
        title_bg = slide.shapes.add_shape(
            1, Inches(0.15), Inches(0), prs.slide_width - Inches(0.15), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = WHITE
        title_bg.line.fill.background()

        # 标题文字
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(10.5), Inches(0.7))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        set_p_font(p, s_title, 28, PRIMARY, bold=True)
        add_entrance_anim(slide, tx, 100, 'fade')

        # 副标题
        if subtitle:
            tx_sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(10), Inches(0.35))
            p_sub = tx_sub.text_frame.paragraphs[0]
            set_p_font(p_sub, subtitle, 15, MID_TEXT)

        # 页码标签
        pg_shape = slide.shapes.add_shape(
            5, Inches(12.2), Inches(0.2), Inches(0.9), Inches(0.5)
        )
        pg_shape.fill.solid()
        pg_shape.fill.fore_color.rgb = PRIMARY
        pg_shape.line.fill.background()
        pg_tf = pg_shape.text_frame
        pg_tf.word_wrap = False
        pg_p = pg_tf.paragraphs[0]
        set_p_font(pg_p, f'{idx+1}/{len(slides_data)}', 12, WHITE, bold=True)
        pg_p.alignment = PP_ALIGN.CENTER

        # 蓝色分隔线
        sep = slide.shapes.add_shape(
            1, Inches(0.6), Inches(1.15), Inches(12), Emu(14000)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = ACCENT
        sep.line.fill.background()

        # ── 内容区域 ──
        content_y = 1.4
        has_code = bool(code.strip())

        if has_code:
            bullet_width = 5.5
            code_left = 6.5
            code_width = 6.2
        else:
            bullet_width = 11.8
            code_left = 0
            code_width = 0

        # 分离表格行和普通要点
        normal_bullets = []
        table_rows = []
        pending_table = []
        if bullets:
            for b in bullets:
                text = str(b)
                if text.startswith('|'):
                    pending_table.append(text)
                else:
                    if pending_table:
                        table_rows.extend(pending_table)
                        pending_table = []
                    normal_bullets.append(text)
            if pending_table:
                table_rows.extend(pending_table)

        # ── 要点列表 — 卡片式 ──
        bullet_count = 0
        if normal_bullets:
            # 创建白色卡片背景
            card = add_card(slide, Inches(0.5), Inches(content_y - 0.1),
                           Inches(bullet_width + 0.3),
                           Inches(5.6 if not table_rows else 2.8))
            add_entrance_anim(slide, card, 200, 'fade')

            tx_b = slide.shapes.add_textbox(
                Inches(0.8), Inches(content_y),
                Inches(bullet_width), Inches(5.2 if not table_rows else 2.5)
            )
            tf_b = tx_b.text_frame
            tf_b.word_wrap = True
            add_entrance_anim(slide, tx_b, 300, 'fade')

            for bi, text in enumerate(normal_bullets):
                if bi == 0:
                    p = tf_b.paragraphs[0]
                else:
                    p = tf_b.add_paragraph()

                if text.startswith('\u3010'):
                    # 【标签】样式 — 使用强调色
                    set_p_font(p, text, 16, SECONDARY, bold=True)
                    p.space_before = Pt(14)
                elif text.startswith('  \u00b7'):
                    set_p_font(p, text, 14, MID_TEXT)
                    p.space_before = Pt(4)
                    p.level = 1
                else:
                    # 使用圆形图标前缀
                    set_p_font(p, f'\u25b8 {text}', 15, DARK_TEXT)
                    p.space_before = Pt(7)
                bullet_count += 1

        # 表格
        if table_rows:
            if normal_bullets:
                tbl_top = content_y + min(bullet_count * 0.32 + 0.4, 3.2)
            else:
                tbl_top = content_y
            tbl_top = min(tbl_top, 4.5)
            add_table(slide, table_rows, 0.8, tbl_top, bullet_width)

        # ── 代码块 — 暗色主题卡片 ──
        if has_code:
            # 代码卡片背景
            code_card = slide.shapes.add_shape(
                5, Inches(code_left - 0.1), Inches(content_y - 0.1),
                Inches(code_width + 0.2), Inches(5.4)
            )
            code_card.fill.solid()
            code_card.fill.fore_color.rgb = CODE_BG
            code_card.line.color.rgb = RGBColor(0x40, 0x40, 0x50)
            code_card.line.width = Pt(1)
            add_entrance_anim(slide, code_card, 400, 'fade')

            # 代码标题栏
            code_title = slide.shapes.add_shape(
                1, Inches(code_left - 0.1), Inches(content_y - 0.1),
                Inches(code_width + 0.2), Inches(0.35)
            )
            code_title.fill.solid()
            code_title.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x3A)
            code_title.line.fill.background()
            ct_tf = code_title.text_frame
            ct_p = ct_tf.paragraphs[0]
            set_p_font(ct_p, '  \u25cf \u25cf \u25cf  Code', 10, RGBColor(0x88, 0x88, 0x99), name='Consolas')

            # 代码内容
            code_box = slide.shapes.add_textbox(
                Inches(code_left), Inches(content_y + 0.3),
                Inches(code_width), Inches(4.9)
            )
            tf_c = code_box.text_frame
            tf_c.word_wrap = True
            code_lines = code.strip().split('\n')
            for ci, cl in enumerate(code_lines):
                if ci == 0:
                    p = tf_c.paragraphs[0]
                else:
                    p = tf_c.add_paragraph()
                set_p_font(p, cl, 12, CODE_FG, name='Consolas')
                p.space_before = Pt(2)
            add_entrance_anim(slide, code_box, 500, 'fade')

        # 备注
        auto_notes = []
        auto_notes.append(s_title)
        if subtitle:
            auto_notes.append(subtitle)
        for b in normal_bullets[:5]:
            auto_notes.append(f'- {b}')
        final_notes = notes if notes else '\n'.join(auto_notes)
        if final_notes:
            slide.notes_slide.notes_text_frame.text = final_notes

    # ═══════════════════════════════════════════════════════════════
    # 结束页 — 渐变 + 大号"谢谢" + Q&A
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, PRIMARY, RGBColor(0x1A, 0x23, 0x7E))
    add_slide_transition(slide, 'fade')

    # 装饰圆
    deco2 = slide.shapes.add_shape(
        9, Inches(8), Inches(4), Inches(8), Inches(8)
    )
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = SECONDARY
    deco2.fill.fore_color.brightness = 0.05
    deco2.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    set_p_font(p, '\u8c22\u8c22\uff01', 54, WHITE, bold=True)
    p.alignment = PP_ALIGN.CENTER
    add_entrance_anim(slide, tx, 200, 'fade')

    p2 = tf.add_paragraph()
    set_p_font(p2, 'Questions & Answers', 22, RGBColor(0xBB, 0xDD, 0xFF))
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    # 底部信息
    tx3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(0.8))
    tf3 = tx3.text_frame
    p3 = tf3.paragraphs[0]
    set_p_font(p3, '\u79fb\u52a8\u5e94\u7528\u5f00\u53d1\u77e5\u8bc6\u56fe\u8c31\u6559\u5b66\u7cfb\u7edf', 14, RGBColor(0x88, 0x99, 0xCC))
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf3.add_paragraph()
    import datetime
    set_p_font(p4, datetime.datetime.now().strftime('%Y\u5e74%m\u6708'), 12, RGBColor(0x77, 0x88, 0xBB))
    p4.alignment = PP_ALIGN.CENTER

    # 保存
    os.makedirs(os.path.dirname(output), exist_ok=True)
    prs.save(output)
    print(f"PPTX saved: {output}")

if __name__ == '__main__':
    main()
